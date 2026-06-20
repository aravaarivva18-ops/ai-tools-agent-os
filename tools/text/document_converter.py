"""Document Converter module for transforming PDF, DOCX, XLSX, PPTX to Markdown."""

import os

import docx
import openpyxl
import pptx
import pypdf


def convert_pdf_to_markdown(filepath: str | os.PathLike[str]) -> str:
    """Extracts text from a PDF file and formats it as Markdown."""
    markdown_lines = []
    # Convert PathLike to string for pypdf
    reader = pypdf.PdfReader(str(filepath))
    total_pages = len(reader.pages)
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        markdown_lines.append(f"## Page {i + 1} of {total_pages}\n")
        markdown_lines.append(text.strip())
        markdown_lines.append("\n")
    return "\n".join(markdown_lines)


def convert_docx_to_markdown(filepath: str | os.PathLike[str]) -> str:
    """Extracts text, headers, and tables from a DOCX file into Markdown."""
    from docx.text.paragraph import Paragraph

    doc = docx.Document(str(filepath))
    markdown_lines = []

    for element in doc.element.body:
        # Check if element is a paragraph
        if element.tag.endswith("p"):
            p = Paragraph(element, doc)
            text = p.text.strip()
            if not text:
                continue

            # Determine paragraph style / heading level
            style_name = (p.style.name or "").lower() if p.style else ""
            if "heading 1" in style_name:
                markdown_lines.append(f"\n# {text}\n")
            elif "heading 2" in style_name:
                markdown_lines.append(f"\n## {text}\n")
            elif "heading 3" in style_name:
                markdown_lines.append(f"\n### {text}\n")
            elif "list" in style_name or p.paragraph_format.left_indent:
                markdown_lines.append(f"* {text}")
            else:
                markdown_lines.append(text)

        # Check if element is a table
        elif element.tag.endswith("tbl"):
            table = docx.table.Table(element, doc)
            markdown_lines.append("")
            for i, row in enumerate(table.rows):
                row_cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                # Filter out adjacent duplicate cells (docx merged cells artifact)
                cleaned_cells: list[str] = []
                for cell in row_cells:
                    if not cleaned_cells or cleaned_cells[-1] != cell:
                        cleaned_cells.append(cell)

                markdown_lines.append("| " + " | ".join(cleaned_cells) + " |")
                if i == 0:
                    markdown_lines.append(
                        "| " + " | ".join(["---"] * len(cleaned_cells)) + " |"
                    )
            markdown_lines.append("")

    return "\n".join(markdown_lines).strip()


def convert_xlsx_to_markdown(filepath: str | os.PathLike[str]) -> str:
    """Converts Excel spreadsheets (all sheets) into Markdown tables."""
    wb = openpyxl.load_workbook(filepath, data_only=True)
    markdown_lines = []

    for name in wb.sheetnames:
        sheet = wb[name]
        markdown_lines.append(f"\n# Sheet: {name}\n")

        # Find maximum dimensions
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            markdown_lines.append("*Empty Sheet*\n")
            continue

        # Clean rows (remove trailing empty rows)
        active_rows = []
        for r in rows:
            if any(val is not None for val in r):
                active_rows.append(r)

        if not active_rows:
            markdown_lines.append("*Empty Sheet*\n")
            continue

        # Convert cells to string representation
        for i, r in enumerate(active_rows):
            cells = [
                str(val).replace("\n", " ").strip() if val is not None else ""
                for val in r
            ]
            markdown_lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                markdown_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")

        markdown_lines.append("")

    return "\n".join(markdown_lines).strip()


def convert_pptx_to_markdown(filepath: str | os.PathLike[str]) -> str:
    """Extracts text content and structures from PowerPoint slides into Markdown."""
    prs = pptx.Presentation(str(filepath))
    markdown_lines = []

    for i, slide in enumerate(prs.slides):
        markdown_lines.append(f"\n# Slide {i + 1}\n")

        # Process title first if available
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                markdown_lines.append(f"## {title_text}\n")

        for shape in slide.shapes:
            # Avoid repeating title shape text
            if shape == slide.shapes.title:
                continue

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        markdown_lines.append(
                            f"* {text}"
                            if len(shape.text_frame.paragraphs) > 1
                            else text
                        )

            # Extract table content inside slides
            elif shape.has_table:
                markdown_lines.append("")
                table = shape.table
                for row_idx in range(len(table.rows)):
                    cells = []
                    for col_idx in range(len(table.columns)):
                        cell_text = (
                            table.cell(row_idx, col_idx).text.replace("\n", " ").strip()
                        )
                        cells.append(cell_text)
                    markdown_lines.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        markdown_lines.append(
                            "| " + " | ".join(["---"] * len(cells)) + " |"
                        )
                markdown_lines.append("")

        # Extract presenter notes if available
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                markdown_lines.append(f"\n> **Notes:** {notes}\n")

    return "\n".join(markdown_lines).strip()


def convert_to_markdown(filepath: str | os.PathLike[str]) -> str:
    """Detects file extension and converts the document to Markdown."""
    _, ext = os.path.splitext(str(filepath).lower())
    if ext == ".pdf":
        return convert_pdf_to_markdown(filepath)
    elif ext == ".docx":
        return convert_docx_to_markdown(filepath)
    elif ext == ".xlsx":
        return convert_xlsx_to_markdown(filepath)
    elif ext == ".pptx":
        return convert_pptx_to_markdown(filepath)
    else:
        raise ValueError(f"Unsupported document format: {ext}")
